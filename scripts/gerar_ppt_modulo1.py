#!/usr/bin/env python3
"""Gera o PPTX do Módulo 1 — SQL Fundamentals (NTT DATA UNIVERSITY)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Cores ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0B, 0x15, 0x35)
ORANGE  = RGBColor(0xF5, 0x81, 0x1F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xAA, 0xAA)
CODE_BG = RGBColor(0x14, 0x25, 0x50)
CODE_FG = RGBColor(0xA8, 0xD8, 0xFF)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=NAVY):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    p.alignment = align
    return tb


def bullets(slide, items, x, y, w, h, size=16, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(size * 0.25)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"


def code_box(slide, code, x, y, w, h, fsize=11):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.width = Pt(0)
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(fsize)
        r.font.color.rgb = CODE_FG
        r.font.name = "Courier New"


def bookmark(slide):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.18), Inches(0.85)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = ORANGE
    s.line.width = Pt(0)


def header(slide):
    txt(slide, "NTT DATA UNIVERSITY",
        x=9.5, y=0.15, w=3.6, h=0.45,
        size=11, bold=True, color=WHITE,
        align=PP_ALIGN.RIGHT)


def footer(slide, n):
    txt(slide, "© 2024 NTT DATA, Inc.",
        x=0.3, y=7.1, w=3, h=0.35, size=9, color=GRAY)
    txt(slide, str(n),
        x=6.4, y=7.1, w=0.8, h=0.35, size=9, color=GRAY,
        align=PP_ALIGN.CENTER)
    txt(slide, "NTT DATA UNIVERSITY",
        x=10.5, y=7.1, w=2.7, h=0.35,
        size=9, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)


def chrome(slide, n):
    bookmark(slide)
    header(slide)
    footer(slide, n)


# ── Slides ───────────────────────────────────────────────────────────────────

def s01_cover(prs):
    s = blank(prs); set_bg(s)
    bookmark(s)
    txt(s, "NTT DATA UNIVERSITY",
        x=9.5, y=0.2, w=3.6, h=0.5,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    txt(s, "SQL Fundamentals",
        x=0.8, y=2.5, w=8, h=1.5,
        size=52, color=WHITE, font="Georgia")
    txt(s, "Módulo 1  |  Trilha de Engenharia de Dados",
        x=0.8, y=4.15, w=8, h=0.7,
        size=20, color=ORANGE)
    footer(s, 1)


def s02_agenda(prs):
    s = blank(prs); set_bg(s); chrome(s, 2)
    txt(s, "Agenda", x=0.8, y=0.9, w=10, h=0.8,
        size=36, color=WHITE, font="Georgia")
    items = [
        "01   Abertura                                    10 min",
        "02   Bloco de Dúvidas                            30 min",
        "03   Exercício em Grupo — Ex 3.3                 60 min",
        "04   Fechamento & Preview Módulo 2               20 min",
    ]
    bullets(s, items, x=0.8, y=2.1, w=11, h=4, size=22)


def s_section(prs, title, sub, n):
    s = blank(prs); set_bg(s)
    bookmark(s); footer(s, n)
    txt(s, title, x=0.8, y=2.7, w=11, h=1.5,
        size=54, color=WHITE, font="Georgia")
    txt(s, sub, x=0.8, y=4.35, w=9, h=0.7,
        size=20, color=ORANGE)


def s04_verificacao(prs):
    s = blank(prs); set_bg(s); chrome(s, 4)
    txt(s, "Verificação de Ambiente",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    txt(s, "Rode a query abaixo e confirme que o resultado é 3000:",
        x=0.8, y=1.65, w=11.5, h=0.45, size=16)
    code_box(s,
        "SELECT COUNT(*) AS total_vendas FROM vendas;\n-- Resultado esperado: 3000",
        x=0.8, y=2.2, w=9, h=0.85, fsize=14)
    txt(s, "Banco não gerado? Execute:",
        x=0.8, y=3.2, w=9, h=0.4, size=14, color=ORANGE)
    code_box(s, "python recursos/setup_db.py",
        x=0.8, y=3.7, w=6, h=0.5, fsize=13)
    txt(s, "Formas de conectar ao banco:",
        x=0.8, y=4.4, w=10, h=0.4, size=14)
    bullets(s, [
        "  SQLite CLI :  sqlite3 recursos/dados.db   →   .headers on   →   .mode column",
        "  Python     :  import sqlite3; conn = sqlite3.connect('recursos/dados.db')",
    ], x=0.8, y=4.9, w=12, h=1.2, size=13)


def s06_where_having(prs):
    s = blank(prs); set_bg(s); chrome(s, 6)
    txt(s, "WHERE vs HAVING",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    txt(s, "Regra: condição sobre agregação → HAVING.  Qualquer outra → WHERE.",
        x=0.8, y=1.6, w=11.5, h=0.45, size=15, color=ORANGE)
    code_box(s,
        "-- ❌ ERRADO\nSELECT cat.nome, SUM(p.preco) AS total\n"
        "FROM produtos p\nINNER JOIN categorias cat\n"
        "  ON p.categoria_id = cat.categoria_id\n"
        "WHERE SUM(p.preco) > 1000  -- erro!\nGROUP BY cat.nome;",
        x=0.8, y=2.15, w=5.8, h=2.1, fsize=11)
    code_box(s,
        "-- ✅ CORRETO\nSELECT cat.nome, SUM(p.preco) AS total\n"
        "FROM produtos p\nINNER JOIN categorias cat\n"
        "  ON p.categoria_id = cat.categoria_id\n"
        "GROUP BY cat.nome\nHAVING SUM(p.preco) > 1000;",
        x=6.8, y=2.15, w=5.8, h=2.1, fsize=11)
    txt(s, "→  Execute ambas no banco para ver os erros e os resultados.",
        x=0.8, y=4.45, w=11, h=0.4, size=13, color=GRAY)


def s07_cte(prs):
    s = blank(prs); set_bg(s); chrome(s, 7)
    txt(s, "CTE vs Subquery",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    txt(s, "Mesma performance na maioria dos casos.  A diferença é legibilidade e manutenção.",
        x=0.8, y=1.6, w=11.5, h=0.45, size=15, color=ORANGE)
    code_box(s,
        "-- Subquery — difícil de ler\nSELECT nome FROM clientes\n"
        "WHERE cliente_id IN (\n"
        "    SELECT cliente_id\n"
        "    FROM vendas\n"
        "    GROUP BY cliente_id\n"
        "    HAVING SUM(valor_total) > 1000\n);",
        x=0.8, y=2.15, w=5.8, h=2.5, fsize=11)
    code_box(s,
        "-- CTE — fácil de ler e testar\nWITH clientes_valiosos AS (\n"
        "    SELECT cliente_id\n"
        "    FROM vendas\n"
        "    GROUP BY cliente_id\n"
        "    HAVING SUM(valor_total) > 1000\n)\n"
        "SELECT c.nome FROM clientes c\nINNER JOIN clientes_valiosos cv\n"
        "    ON c.cliente_id = cv.cliente_id;",
        x=6.8, y=2.15, w=5.8, h=2.5, fsize=11)
    txt(s, "Orientação: prefira CTE em qualquer query com mais de uma camada de lógica.",
        x=0.8, y=4.85, w=11, h=0.4, size=13, color=GRAY)


def s08_left_join(prs):
    s = blank(prs); set_bg(s); chrome(s, 8)
    txt(s, "LEFT JOIN — Encontrar o que não existe",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    txt(s, "Padrão: LEFT JOIN  +  WHERE <pk_tabela_direita> IS NULL",
        x=0.8, y=1.6, w=11.5, h=0.45, size=15, color=ORANGE)
    code_box(s,
        "-- Clientes que nunca compraram\nSELECT c.nome\n"
        "FROM clientes AS c\nLEFT JOIN vendas AS v\n"
        "    ON c.cliente_id = v.cliente_id\n"
        "WHERE v.venda_id IS NULL;    -- checar PK da tabela direita",
        x=0.8, y=2.2, w=9, h=1.8, fsize=13)
    txt(s, "⚠  Armadilha:",
        x=0.8, y=4.2, w=2.2, h=0.45, size=15, bold=True, color=ORANGE)
    txt(s,
        "Prefira checar a chave primária (venda_id), não a FK (cliente_id).\n"
        "Se cliente_id for NOT NULL em vendas, o comportamento pode ser inesperado.",
        x=0.8, y=4.7, w=11, h=0.8, size=14, color=GRAY)


def s09_datas(prs):
    s = blank(prs); set_bg(s); chrome(s, 9)
    txt(s, "Datas no SQLite — strftime",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    txt(s, "SQLite não tem DATE_TRUNC. Use strftime('%Y-%m', coluna) para agrupar por mês.",
        x=0.8, y=1.6, w=11.5, h=0.45, size=15, color=ORANGE)
    code_box(s,
        "-- Receita por mês\nSELECT\n"
        "    strftime('%Y-%m', data_venda) AS mes,\n"
        "    SUM(valor_total)              AS receita_mes\nFROM vendas\n"
        "GROUP BY strftime('%Y-%m', data_venda)\nORDER BY mes;\n\n"
        "-- Filtrar apenas 2024\nWHERE strftime('%Y', data_venda) = '2024'",
        x=0.8, y=2.2, w=9, h=3.0, fsize=13)


def s11_query_original(prs):
    s = blank(prs); set_bg(s); chrome(s, 11)
    txt(s, "Ex 3.3 — A Query Original",
        x=0.8, y=0.7, w=11, h=0.65,
        size=26, bold=True, color=WHITE, font="Georgia")
    txt(s, '"O que essa query faz? Explique em uma frase.  |  Onde você vê desperdício?"',
        x=0.8, y=1.4, w=11.5, h=0.45, size=14, color=ORANGE)
    code_box(s,
        "SELECT *\nFROM (\n    SELECT * FROM clientes\n) AS todos_clientes\n"
        "WHERE cliente_id IN (\n    SELECT cliente_id\n    FROM vendas\n"
        "    WHERE valor_total > (\n        SELECT AVG(valor_total)\n"
        "        FROM vendas\n        WHERE cliente_id IN (\n"
        "            SELECT cliente_id FROM clientes\n        )\n    )\n);",
        x=0.8, y=2.0, w=7.5, h=4.7, fsize=12)


def s12_problemas(prs):
    s = blank(prs); set_bg(s); chrome(s, 12)
    txt(s, "Onde estão os problemas?",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    bullets(s, [
        "1.  SELECT *  — precisamos de todas as colunas? Quais o resultado final precisa?",
        "",
        "2.  FROM (SELECT * FROM clientes) — o que essa subquery adiciona?",
        "     Poderíamos usar FROM clientes diretamente?",
        "",
        "3.  WHERE cliente_id IN (SELECT cliente_id FROM clientes) — todo cliente_id",
        "     em vendas já existe em clientes (FK). Essa subquery filtra alguma coisa?",
        "",
        "4.  Subquery escalar no WHERE — o banco calcula AVG(valor_total)",
        "     uma vez total, ou uma vez para cada linha de vendas?",
    ], x=0.8, y=1.8, w=12, h=5, size=15)


def s13_refatoracao(prs):
    s = blank(prs); set_bg(s); chrome(s, 13)
    txt(s, "Query Refatorada — Construída Coletivamente",
        x=0.8, y=0.7, w=11, h=0.65,
        size=26, bold=True, color=WHITE, font="Georgia")
    code_box(s,
        "WITH media_vendas AS (\n"
        "    SELECT AVG(valor_total) AS media_geral\n"
        "    FROM vendas                            -- passo 1: isolar a média\n"
        "),\n"
        "clientes_acima_da_media AS (\n"
        "    SELECT DISTINCT v.cliente_id\n"
        "    FROM vendas AS v, media_vendas AS m\n"
        "    WHERE v.valor_total > m.media_geral    -- passo 2: filtrar\n"
        ")\n"
        "SELECT                                     -- passo 3: dados do cliente\n"
        "    c.cliente_id, c.nome, c.email,\n"
        "    c.cidade, c.estado, c.data_cadastro\n"
        "FROM clientes AS c\n"
        "INNER JOIN clientes_acima_da_media AS cam\n"
        "    ON c.cliente_id = cam.cliente_id\n"
        "ORDER BY c.nome;",
        x=0.8, y=1.5, w=11.5, h=5.5, fsize=12)


def s14_explain(prs):
    s = blank(prs); set_bg(s); chrome(s, 14)
    txt(s, "EXPLAIN QUERY PLAN — Analisando o Plano",
        x=0.8, y=0.8, w=11, h=0.7,
        size=26, bold=True, color=WHITE, font="Georgia")
    code_box(s, "EXPLAIN QUERY PLAN\n<cole a query aqui>;",
        x=0.8, y=1.65, w=8, h=0.7, fsize=13)
    bullets(s, [
        "SCAN   →  varredura completa da tabela (pode ser lento em produção)",
        "SEARCH →  uso de índice (mais eficiente)",
        "",
        "O otimizador já eliminou alguma redundância automaticamente?",
        "O que muda criando um índice em vendas(valor_total)?",
        "",
        "Nota: com 3.000 linhas a diferença é imperceptível.",
        "O exercício cria o hábito — o ganho aparece em tabelas com milhões de linhas.",
    ], x=0.8, y=2.55, w=12, h=4, size=15)


def s16_objetivos(prs):
    s = blank(prs); set_bg(s); chrome(s, 16)
    txt(s, "O que cobrimos hoje",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    rows = [
        ("✅  SELECT, WHERE, ORDER BY",     "Seção 1 + Ex 1.1, 1.2, 1.3"),
        ("✅  JOINs (INNER, LEFT, RIGHT)",  "Seção 2 + Ex 2.1, 2.2"),
        ("✅  GROUP BY e Agregações",        "Seção 3 + Ex 2.3"),
        ("✅  CTEs",                          "Seção 4 + Ex 3.2, 3.3"),
        ("✅  Window Functions",              "Seção 5 + Ex 3.1, 3.2"),
        ("✅  Performance Básica",            "Seção 6 + Ex 3.3"),
    ]
    for i, (obj, cob) in enumerate(rows):
        y = 1.85 + i * 0.72
        txt(s, obj,  x=0.8, y=y, w=7.5, h=0.6, size=16)
        txt(s, cob,  x=8.5, y=y, w=4.5, h=0.6, size=14, color=ORANGE)


def s17_preview(prs):
    s = blank(prs); set_bg(s); chrome(s, 17)
    txt(s, "Próximo: Módulo 2 — Modelagem de Dados",
        x=0.8, y=0.8, w=11, h=0.7,
        size=28, bold=True, color=WHITE, font="Georgia")
    bullets(s, [
        "▸  Modelagem relacional — entidades, atributos, cardinalidade",
        "▸  Normalização — 1FN, 2FN, 3FN e quando desnormalizar",
        "▸  Modelagem dimensional — tabelas fato e dimensão",
        "▸  Star Schema vs Snowflake Schema",
        "▸  Como o schema do dados.db se encaixa em cada modelo",
    ], x=0.8, y=1.9, w=11.5, h=3.5, size=19)
    txt(s,
        "Desafio antes da próxima sessão: revisar Ex 3.2 e experimentar "
        "ROWS BETWEEN vs RANGE BETWEEN na window function.",
        x=0.8, y=5.9, w=11.5, h=0.6, size=13, color=ORANGE)


def s18_closing(prs):
    s = blank(prs); set_bg(s)
    bookmark(s)
    txt(s, "NTT DATA UNIVERSITY",
        x=2.5, y=3.2, w=8.5, h=0.9,
        size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    footer(s, 18)


def build():
    prs = new_prs()
    s01_cover(prs)
    s02_agenda(prs)
    s_section(prs, "Abertura", "10 minutos", 3)
    s04_verificacao(prs)
    s_section(prs, "Bloco de Dúvidas", "30 minutos", 5)
    s06_where_having(prs)
    s07_cte(prs)
    s08_left_join(prs)
    s09_datas(prs)
    s_section(prs, "Exercício em Grupo", "60 minutos — Ex 3.3", 10)
    s11_query_original(prs)
    s12_problemas(prs)
    s13_refatoracao(prs)
    s14_explain(prs)
    s_section(prs, "Fechamento", "20 minutos", 15)
    s16_objetivos(prs)
    s17_preview(prs)
    s18_closing(prs)

    out = "modulos/01-sql-fundamentals/sessao-ao-vivo-modulo1.pptx"
    prs.save(out)
    print(f"Gerado: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
